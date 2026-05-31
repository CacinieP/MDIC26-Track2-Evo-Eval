"""
MinerU Parser Tool - Core document parsing via MinerU v1.3.x.

Wraps MinerU's magic-pdf pipeline for:
- Layout detection & reading order
- Text extraction with OCR fallback
- Table structure recognition
- Image/figure extraction
- Formula detection

Supports PDF, images, DOCX, PPTX, and HTML files.
"""

from __future__ import annotations

import asyncio
import base64
import io
import json
import os
import shutil
import time
from pathlib import Path
from typing import Any

import cv2
import numpy as np
from loguru import logger
from PIL import Image

# ---------------------------------------------------------------------------
# Supported file-type sets
# ---------------------------------------------------------------------------
_PDF_EXTS = {".pdf"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}
_OFFICE_EXTS = {".docx", ".pptx"}
_HTML_EXTS = {".html", ".htm"}
_SUPPORTED_EXTS = _PDF_EXTS | _IMAGE_EXTS | _OFFICE_EXTS | _HTML_EXTS


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _check_mineru_available() -> bool:
    """Return True if MinerU (magic_pdf) is importable."""
    try:
        import magic_pdf  # noqa: F401
        return True
    except ImportError:
        return False


def _suffix(path: str | Path) -> str:
    return Path(path).suffix.lower()


# ---------------------------------------------------------------------------
# Image pre-processing utilities (applied before OCR)
# ---------------------------------------------------------------------------

class ImagePreprocessor:
    """Pre-process images for better OCR quality."""

    @staticmethod
    def enhance(image: np.ndarray) -> np.ndarray:
        """Run a lightweight enhancement pipeline:
        1. Convert to grayscale if colour
        2. CLAHE contrast normalisation
        3. Adaptive denoising
        4. Deskew (simple projection-based)
        """
        gray = image if len(image.shape) == 2 else cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

        # CLAHE
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(gray)

        # Denoise
        enhanced = cv2.fastNlMeansDenoising(enhanced, h=10, templateWindowSize=7, searchWindowSize=21)

        # Deskew
        enhanced = ImagePreprocessor._deskew(enhanced)

        return enhanced

    @staticmethod
    def _deskew(image: np.ndarray) -> np.ndarray:
        """Correct skew using min-area bounding rect of all foreground pixels."""
        thresh = cv2.threshold(image, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)[1]
        coords = np.column_stack(np.where(thresh > 0))
        if coords.shape[0] < 50:
            return image
        angle = cv2.minAreaRect(coords)[-1]
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
        if abs(angle) < 0.5:
            return image
        (h, w) = image.shape[:2]
        center = (w // 2, h // 2)
        matrix = cv2.getRotationMatrix2D(center, angle, 1.0)
        rotated = cv2.warpAffine(image, matrix, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
        return rotated


# ---------------------------------------------------------------------------
# Main parser class
# ---------------------------------------------------------------------------

class MinerUParser:
    """
    MinerU document parser tool.

    Usage::

        parser = MinerUParser(config)
        result = await parser.execute({"file_path": "/path/to/doc.pdf"}, context)

    The *config* dict accepts the following keys (all optional):

    - ``model_dir``  – directory for MinerU models (default ``"./models"``)
    - ``device``     – ``"cuda"`` or ``"cpu"``  (default ``"cuda"``)
    - ``output_dir`` – base directory for MinerU intermediate outputs
    - ``preprocess`` – ``bool``, run image enhancement before OCR (default ``True``)
    - ``table_enable``   – enable table recognition (default ``True``)
    - ``formula_enable`` – enable formula recognition (default ``True``)
    - ``ocr_lang``       – OCR language hint, e.g. ``"ch"``, ``"en"`` (default ``"auto"``)
    """

    def __init__(self, config: dict | None = None):
        self.config = config or {}
        self.model_dir = self.config.get("model_dir", "./models")
        self.device = self.config.get("device", "cuda")
        self.output_base = Path(self.config.get("output_dir", "./data/output"))
        self.preprocess_enabled = self.config.get("preprocess", True)
        self.table_enable = self.config.get("table_enable", True)
        self.formula_enable = self.config.get("formula_enable", True)
        self.ocr_lang = self.config.get("ocr_lang", "auto")
        self._mineru_available = _check_mineru_available()
        self._ocr_engine: Any | None = None  # Lazy-initialised PaddleOCR singleton

        # Cloud API configuration
        raw_token = self.config.get("api_token", os.environ.get("MINERU_API_TOKEN", ""))
        # Filter out unresolved env var placeholders like "${MINERU_API_TOKEN}"
        self.api_token = raw_token if raw_token and not raw_token.startswith("${") else ""
        self.api_mode = self.config.get("api_mode", "auto")  # auto | cloud | local

        if self._mineru_available:
            logger.info("MinerUParser: MinerU (magic_pdf) detected and ready")
        else:
            logger.warning("MinerUParser: MinerU local models not available — will use cloud API or fallback")

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    # ------------------------------------------------------------------
    # OCR engine (lazy singleton)
    # ------------------------------------------------------------------

    def _get_ocr_engine(self) -> Any:
        """Return a lazily-initialised PaddleOCR engine (singleton per parser)."""
        if self._ocr_engine is None:
            from paddleocr import PaddleOCR
            self._ocr_engine = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
            logger.info("MinerUParser: PaddleOCR engine initialised (cached)")
        return self._ocr_engine

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def execute(self, params: dict[str, Any], context: dict | None = None) -> dict:
        """
        Parse a document.

        Args:
            params:  Must contain ``file_path`` (str | Path).
                     Optional keys:
                     - ``output_format``: ``"markdown"`` | ``"blocks"`` (default ``"markdown"``)
                     - ``preprocess``: override instance-level preprocess flag
            context: Results carried from previous pipeline steps (unused by this tool).

        Returns:
            A structured dict::

                {
                    "source_file": str,
                    "pages": int,
                    "content_list": list[dict],
                    "markdown": str,
                    "tables": list[dict],
                    "images": list[dict],
                    "metadata": dict,
                }
        """
        file_path = params.get("file_path")
        if not file_path:
            # Demo / dry-run mode: return empty result instead of crashing
            logger.warning("MinerUParser: no file_path provided, returning empty result")
            return {
                "source_file": "",
                "pages": 0,
                "content_list": [],
                "markdown": "",
                "tables": [],
                "images": [],
                "metadata": {"parser": "mineru", "mode": "demo"},
            }

        file_path = Path(file_path)

        suffix = _suffix(file_path)
        if suffix not in _SUPPORTED_EXTS:
            raise ValueError(
                f"Unsupported file format '{suffix}'. "
                f"Supported: {sorted(_SUPPORTED_EXTS)}"
            )

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        preprocess = params.get("preprocess", self.preprocess_enabled)
        t0 = time.perf_counter()
        logger.info(f"MinerUParser.execute | file={file_path.name}  suffix={suffix}  preprocess={preprocess}")

        # Route to the appropriate handler
        # When api_mode is "cloud", skip local processing entirely
        if self.api_mode == "cloud":
            logger.info(f"api_mode=cloud, routing directly to cloud API for {file_path.name}")
            cloud_result = await self._parse_via_cloud_api(file_path)
            if cloud_result is not None:
                result = cloud_result
            else:
                logger.info("Cloud API unavailable — falling back to basic parsers")
                result = await self._fallback_parse(file_path)
        else:
            try:
                if suffix in _PDF_EXTS:
                    result = await self._handle_pdf(file_path, preprocess)
                elif suffix in _IMAGE_EXTS:
                    result = await self._handle_image(file_path, preprocess)
                elif suffix in _OFFICE_EXTS:
                    result = await self._handle_office(file_path)
                elif suffix in _HTML_EXTS:
                    result = await self._handle_html(file_path)
                else:
                    raise ValueError(f"Unhandled suffix: {suffix}")
            except Exception as exc:
                logger.error(f"MinerUParser primary handler failed for {file_path.name}: {exc}")

                # Try MinerU cloud API before falling back to basic parsers
                cloud_result = await self._parse_via_cloud_api(file_path)
                if cloud_result is not None:
                    result = cloud_result
                else:
                    logger.info("Cloud API unavailable or failed — falling back to basic parsers")
                    result = await self._fallback_parse(file_path)

        elapsed = time.perf_counter() - t0
        result["metadata"]["parse_time_s"] = round(elapsed, 3)
        logger.info(
            f"MinerUParser done | file={file_path.name}  pages={result.get('pages', 0)}  "
            f"tables={len(result.get('tables', []))}  images={len(result.get('images', []))}  "
            f"time={elapsed:.2f}s"
        )
        return result

    # ------------------------------------------------------------------
    # PDF handling (MinerU primary path)
    # ------------------------------------------------------------------

    async def _handle_pdf(self, file_path: Path, preprocess: bool) -> dict:
        """Parse a PDF via MinerU's ``PymuDocDataset`` pipeline."""
        if not self._mineru_available:
            return await self._parse_pdf_basic(file_path, preprocess)

        # --- lazy imports so the module loads even when MinerU is absent ---
        from magic_pdf.data.data_reader_writer import FileBasedDataWriter
        from magic_pdf.data.dataset import PymuDocDataset
        from magic_pdf.config.enums import SupportedPdfParseMethod
        from magic_pdf.config.make_content_config import MakeMode, DropMode
        from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze

        logger.info(f"Reading PDF bytes: {file_path.name}")
        pdf_bytes = file_path.read_bytes()

        # Prepare output directories per document
        out_dir = self._make_output_dir(file_path)
        image_dir = out_dir / "images"
        image_dir.mkdir(parents=True, exist_ok=True)

        # Create dataset
        ds = PymuDocDataset(bits=pdf_bytes, lang=self.ocr_lang)
        logger.info(f"PymuDocDataset created | pages={len(ds)}")

        # Classify → OCR or TXT mode
        classify_result = ds.classify()
        ocr_mode = classify_result == SupportedPdfParseMethod.OCR
        logger.info(f"Classification result: {'OCR' if ocr_mode else 'TXT'} mode")

        # If preprocessing is requested and the doc is OCR-mode, enhance page images
        if preprocess and ocr_mode:
            logger.info("Running image pre-processing for OCR-mode PDF")
            pdf_bytes = self._preprocess_pdf_pages(pdf_bytes)
            ds = PymuDocDataset(bits=pdf_bytes, lang=self.ocr_lang)

        # Run inference
        logger.info("Running doc_analyze inference  table_enable={}  formula_enable={}".format(
            self.table_enable, self.formula_enable))
        infer_result = ds.apply(
            doc_analyze,
            ocr=ocr_mode,
            table_enable=self.table_enable,
            formula_enable=self.formula_enable,
        )

        # Create writer for images extracted during pipe
        image_writer = FileBasedDataWriter(image_dir)

        # Run pipe
        if ocr_mode:
            pipe_result = infer_result.pipe_ocr_mode(image_writer)
        else:
            pipe_result = infer_result.pipe_txt_mode(image_writer)

        logger.info("Pipe processing complete")

        # Extract structured content
        content_list: list[dict] = pipe_result.get_content_list(str(image_dir))
        markdown: str = pipe_result.get_markdown(
            str(image_dir),
            md_make_mode=MakeMode.MM_MD,
        )

        # Build structured result
        tables = [item for item in content_list if item.get("type") == "table"]
        images = [item for item in content_list if item.get("type") == "image"]
        equations = [item for item in content_list if item.get("type") == "equation"]
        text_blocks = [item for item in content_list if item.get("type") == "text"]

        return {
            "source_file": str(file_path),
            "pages": len(ds),
            "content_list": content_list,
            "markdown": markdown,
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_pdf",
                "ocr_mode": ocr_mode,
                "total_content_items": len(content_list),
                "text_blocks": len(text_blocks),
                "table_count": len(tables),
                "image_count": len(images),
                "equation_count": len(equations),
                "output_dir": str(out_dir),
            },
        }

    # ------------------------------------------------------------------
    # Image handling
    # ------------------------------------------------------------------

    async def _handle_image(self, file_path: Path, preprocess: bool) -> dict:
        """Parse a single image file via MinerU ``ImageDataset`` or fallback OCR."""
        if self._mineru_available:
            try:
                return await self._parse_image_mineru(file_path, preprocess)
            except Exception as exc:
                logger.warning(f"MinerU image parse failed, falling back: {exc}")

        return await self._parse_image_basic(file_path, preprocess)

    async def _parse_image_mineru(self, file_path: Path, preprocess: bool) -> dict:
        from magic_pdf.data.data_reader_writer import FileBasedDataWriter
        from magic_pdf.data.dataset import ImageDataset
        from magic_pdf.config.make_content_config import MakeMode
        from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze

        img_bytes = file_path.read_bytes()

        # Optional pre-processing
        if preprocess:
            img_bytes = self._preprocess_image_bytes(img_bytes)

        out_dir = self._make_output_dir(file_path)
        image_dir = out_dir / "images"
        image_dir.mkdir(parents=True, exist_ok=True)

        ds = ImageDataset(bits=img_bytes, lang=self.ocr_lang)
        logger.info("ImageDataset created for {}", file_path.name)

        infer_result = ds.apply(
            doc_analyze,
            ocr=True,
            table_enable=self.table_enable,
            formula_enable=self.formula_enable,
        )

        image_writer = FileBasedDataWriter(image_dir)
        pipe_result = infer_result.pipe_ocr_mode(image_writer)

        content_list: list[dict] = pipe_result.get_content_list(str(image_dir))
        markdown: str = pipe_result.get_markdown(
            str(image_dir),
            md_make_mode=MakeMode.MM_MD,
        )

        tables = [item for item in content_list if item.get("type") == "table"]
        images = [item for item in content_list if item.get("type") == "image"]

        return {
            "source_file": str(file_path),
            "pages": 1,
            "content_list": content_list,
            "markdown": markdown,
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_image",
                "ocr_mode": True,
                "output_dir": str(out_dir),
            },
        }

    async def _parse_image_basic(self, file_path: Path, preprocess: bool) -> dict:
        """Fallback: basic OCR on a single image via PaddleOCR."""
        img_bytes = file_path.read_bytes()
        if preprocess:
            img_bytes = self._preprocess_image_bytes(img_bytes)

        text = ""
        try:
            ocr_engine = self._get_ocr_engine()
            nparr = np.frombuffer(img_bytes, np.uint8)
            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
            result = ocr_engine.ocr(img, cls=True)
            if result and result[0]:
                lines = [line[1][0] for line in result[0] if line[1]]
                text = "\n".join(lines)
        except ImportError:
            logger.warning("PaddleOCR not available; skipping OCR for image")
        except Exception as exc:
            logger.warning(f"OCR failed for {file_path.name}: {exc}")

        return {
            "source_file": str(file_path),
            "pages": 1,
            "content_list": [{"type": "text", "text": text}] if text else [],
            "markdown": text,
            "tables": [],
            "images": [],
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "fallback_image_ocr",
            },
        }

    # ------------------------------------------------------------------
    # Office handling (DOCX / PPTX)
    # ------------------------------------------------------------------

    async def _handle_office(self, file_path: Path) -> dict:
        suffix = _suffix(file_path)
        if suffix == ".docx":
            return await self._parse_docx(file_path)
        elif suffix == ".pptx":
            return await self._parse_pptx(file_path)
        raise ValueError(f"Unexpected office suffix: {suffix}")

    async def _parse_docx(self, file_path: Path) -> dict:
        """Parse DOCX — try MinerU ``read_local_office`` first, then python-docx."""
        if self._mineru_available:
            try:
                return await self._parse_office_mineru(file_path)
            except Exception as exc:
                logger.warning(f"MinerU office parse failed, using python-docx: {exc}")

        return await self._parse_docx_basic(file_path)

    async def _parse_pptx(self, file_path: Path) -> dict:
        """Parse PPTX — try MinerU first, then python-pptx."""
        if self._mineru_available:
            try:
                return await self._parse_office_mineru(file_path)
            except Exception as exc:
                logger.warning(f"MinerU office parse failed, using python-pptx: {exc}")

        return await self._parse_pptx_basic(file_path)

    async def _parse_office_mineru(self, file_path: Path) -> dict:
        """Use MinerU's ``read_local_office`` to convert office docs to PDF-equivalent."""
        from magic_pdf.data.read_api import read_local_office
        from magic_pdf.data.data_reader_writer import FileBasedDataWriter
        from magic_pdf.data.dataset import PymuDocDataset
        from magic_pdf.config.make_content_config import MakeMode
        from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze

        logger.info("Converting office doc via MinerU read_local_office: {}", file_path.name)

        # read_local_office returns a list of (pdf_bytes, lang) tuples
        datasets = read_local_office(str(file_path), lang=self.ocr_lang)
        if not datasets:
            raise RuntimeError("read_local_office returned empty result")

        out_dir = self._make_output_dir(file_path)
        image_dir = out_dir / "images"
        image_dir.mkdir(parents=True, exist_ok=True)

        all_content: list[dict] = []
        all_markdown_parts: list[str] = []
        total_pages = 0
        all_tables: list[dict] = []
        all_images: list[dict] = []

        for ds_entry in datasets:
            pdf_bytes = ds_entry  # bytes
            ds = PymuDocDataset(bits=pdf_bytes, lang=self.ocr_lang)
            total_pages += len(ds)

            classify_result = ds.classify()
            ocr_mode = hasattr(classify_result, "value") and "ocr" in str(classify_result).lower()

            infer_result = ds.apply(
                doc_analyze,
                ocr=ocr_mode,
                table_enable=self.table_enable,
                formula_enable=self.formula_enable,
            )

            image_writer = FileBasedDataWriter(image_dir)
            if ocr_mode:
                pipe_result = infer_result.pipe_ocr_mode(image_writer)
            else:
                pipe_result = infer_result.pipe_txt_mode(image_writer)

            content_list = pipe_result.get_content_list(str(image_dir))
            md = pipe_result.get_markdown(str(image_dir), md_make_mode=MakeMode.MM_MD)

            all_content.extend(content_list)
            all_markdown_parts.append(md)
            all_tables.extend([item for item in content_list if item.get("type") == "table"])
            all_images.extend([item for item in content_list if item.get("type") == "image"])

        return {
            "source_file": str(file_path),
            "pages": total_pages,
            "content_list": all_content,
            "markdown": "\n\n---\n\n".join(all_markdown_parts),
            "tables": all_tables,
            "images": all_images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_office",
                "output_dir": str(out_dir),
            },
        }

    async def _parse_docx_basic(self, file_path: Path) -> dict:
        """Fallback DOCX parsing via python-docx."""
        content_list: list[dict] = []
        tables: list[dict] = []
        images: list[dict] = []
        md_parts: list[str] = []
        page_count = 1

        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT

            doc = Document(str(file_path))
            logger.info("python-docx loaded: {} paragraphs, {} tables", len(doc.paragraphs), len(doc.tables))

            out_dir = self._make_output_dir(file_path)
            img_dir = out_dir / "images"
            img_dir.mkdir(parents=True, exist_ok=True)

            # Paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style_name = (para.style.name or "").lower()
                if "heading 1" in style_name:
                    md_parts.append(f"# {text}")
                elif "heading 2" in style_name:
                    md_parts.append(f"## {text}")
                elif "heading 3" in style_name:
                    md_parts.append(f"### {text}")
                else:
                    md_parts.append(text)
                content_list.append({"type": "text", "text": text})

            # Tables
            for idx, table in enumerate(doc.tables):
                rows_data: list[list[str]] = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_cells)
                if rows_data:
                    header = rows_data[0]
                    table_md = self._table_to_markdown(rows_data)
                    md_parts.append(table_md)
                    table_entry = {
                        "type": "table",
                        "table_index": idx,
                        "data": rows_data,
                        "header": header,
                        "rows": len(rows_data),
                        "cols": len(header),
                        "markdown": table_md,
                    }
                    tables.append(table_entry)
                    content_list.append(table_entry)

            # Inline images
            img_idx = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_bytes = rel.target_part.blob
                        img_path = img_dir / f"image_{img_idx}.png"
                        img_path.write_bytes(img_bytes)
                        images.append({
                            "type": "image",
                            "path": str(img_path),
                            "index": img_idx,
                        })
                        img_idx += 1
                    except Exception as exc:
                        logger.warning(f"Failed to extract DOCX image: {exc}")

        except ImportError:
            logger.warning("python-docx not installed; DOCX parsing skipped")
        except Exception as exc:
            logger.error(f"DOCX parsing error: {exc}")

        return {
            "source_file": str(file_path),
            "pages": page_count,
            "content_list": content_list,
            "markdown": "\n\n".join(md_parts),
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "fallback_docx",
            },
        }

    async def _parse_pptx_basic(self, file_path: Path) -> dict:
        """Fallback PPTX parsing via python-pptx."""
        content_list: list[dict] = []
        tables: list[dict] = []
        images: list[dict] = []
        md_parts: list[str] = []
        slide_count = 0

        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slide_count = len(prs.slides)
            logger.info("python-pptx loaded: {} slides", slide_count)

            out_dir = self._make_output_dir(file_path)
            img_dir = out_dir / "images"
            img_dir.mkdir(parents=True, exist_ok=True)

            for slide_idx, slide in enumerate(prs.slides):
                md_parts.append(f"## Slide {slide_idx + 1}")
                content_list.append({"type": "text", "text": f"--- Slide {slide_idx + 1} ---"})

                for shape in slide.shapes:
                    # Text
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                md_parts.append(text)
                                content_list.append({"type": "text", "text": text})

                    # Tables
                    if shape.has_table:
                        tbl = shape.table
                        rows_data: list[list[str]] = []
                        for row in tbl.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            rows_data.append(row_cells)
                        if rows_data:
                            table_md = self._table_to_markdown(rows_data)
                            md_parts.append(table_md)
                            t_entry = {
                                "type": "table",
                                "table_index": len(tables),
                                "slide": slide_idx,
                                "data": rows_data,
                                "rows": len(rows_data),
                                "cols": len(rows_data[0]) if rows_data else 0,
                                "markdown": table_md,
                            }
                            tables.append(t_entry)
                            content_list.append(t_entry)

                    # Images
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            img_bytes = shape.image.blob
                            ext = shape.image.content_type.split("/")[-1]
                            img_path = img_dir / f"slide{slide_idx}_img{len(images)}.{ext}"
                            img_path.write_bytes(img_bytes)
                            images.append({
                                "type": "image",
                                "path": str(img_path),
                                "slide": slide_idx,
                                "index": len(images),
                            })
                        except Exception as exc:
                            logger.warning(f"Failed to extract PPTX image: {exc}")

        except ImportError:
            logger.warning("python-pptx not installed; PPTX parsing skipped")
        except Exception as exc:
            logger.error(f"PPTX parsing error: {exc}")

        return {
            "source_file": str(file_path),
            "pages": slide_count,
            "content_list": content_list,
            "markdown": "\n\n".join(md_parts),
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "fallback_pptx",
                "slides": slide_count,
            },
        }

    # ------------------------------------------------------------------
    # HTML handling
    # ------------------------------------------------------------------

    async def _handle_html(self, file_path: Path) -> dict:
        return await self._parse_html(file_path)

    async def _parse_html(self, file_path: Path) -> dict:
        """Parse HTML via BeautifulSoup, extracting text, tables, and images."""
        content_list: list[dict] = []
        tables: list[dict] = []
        images: list[dict] = []
        md_parts: list[str] = []

        try:
            from bs4 import BeautifulSoup

            html_text = file_path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(html_text, "lxml")

            out_dir = self._make_output_dir(file_path)
            img_dir = out_dir / "images"
            img_dir.mkdir(parents=True, exist_ok=True)

            # Remove script/style
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()

            # Title
            title_tag = soup.find("title")
            if title_tag and title_tag.string:
                md_parts.append(f"# {title_tag.string.strip()}")
                content_list.append({"type": "text", "text": title_tag.string.strip()})

            # Headings & paragraphs
            for tag in soup.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
                text = tag.get_text(strip=True)
                if not text:
                    continue
                if tag.name == "h1":
                    md_parts.append(f"# {text}")
                elif tag.name == "h2":
                    md_parts.append(f"## {text}")
                elif tag.name == "h3":
                    md_parts.append(f"### {text}")
                elif tag.name == "li":
                    md_parts.append(f"- {text}")
                else:
                    md_parts.append(text)
                content_list.append({"type": "text", "text": text})

            # Tables
            for idx, table_tag in enumerate(soup.find_all("table")):
                rows_data: list[list[str]] = []
                for tr in table_tag.find_all("tr"):
                    cells = tr.find_all(["td", "th"])
                    row = [cell.get_text(strip=True) for cell in cells]
                    if row:
                        rows_data.append(row)
                if rows_data:
                    table_md = self._table_to_markdown(rows_data)
                    md_parts.append(table_md)
                    t_entry = {
                        "type": "table",
                        "table_index": idx,
                        "data": rows_data,
                        "rows": len(rows_data),
                        "cols": max(len(r) for r in rows_data),
                        "markdown": table_md,
                    }
                    tables.append(t_entry)
                    content_list.append(t_entry)

            # Images (save base64 or linked)
            for idx, img_tag in enumerate(soup.find_all("img")):
                src = img_tag.get("src", "")
                alt = img_tag.get("alt", "")
                if src.startswith("data:"):
                    # base64 embedded
                    try:
                        header, b64data = src.split(",", 1)
                        ext = header.split("/")[1].split(";")[0]
                        img_bytes = base64.b64decode(b64data)
                        img_path = img_dir / f"html_img_{idx}.{ext}"
                        img_path.write_bytes(img_bytes)
                        images.append({
                            "type": "image",
                            "path": str(img_path),
                            "alt": alt,
                            "index": idx,
                        })
                    except Exception as exc:
                        logger.warning(f"Failed to decode base64 image: {exc}")
                else:
                    images.append({
                        "type": "image",
                        "src": src,
                        "alt": alt,
                        "index": idx,
                    })

        except ImportError:
            logger.warning("BeautifulSoup4 / lxml not installed; HTML parsing skipped")
        except Exception as exc:
            logger.error(f"HTML parsing error: {exc}")

        return {
            "source_file": str(file_path),
            "pages": 1,
            "content_list": content_list,
            "markdown": "\n\n".join(md_parts),
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "fallback_html",
            },
        }

    # ------------------------------------------------------------------
    # PDF fallback (basic)
    # ------------------------------------------------------------------

    async def _parse_pdf_basic(self, file_path: Path, preprocess: bool) -> dict:
        """Basic PDF fallback: extract text via PyMuPDF or pdf2image + OCR."""
        content_list: list[dict] = []
        tables: list[dict] = []
        images: list[dict] = []
        md_parts: list[str] = []
        page_count = 0

        # Attempt PyMuPDF (fitz)
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(file_path))
            page_count = len(doc)
            logger.info(f"PyMuPDF loaded: {page_count} pages")

            out_dir = self._make_output_dir(file_path)
            img_dir = out_dir / "images"
            img_dir.mkdir(parents=True, exist_ok=True)

            for page_num in range(page_count):
                page = doc[page_num]
                text = page.get_text("text").strip()
                if text:
                    md_parts.append(text)
                    content_list.append({"type": "text", "text": text, "page": page_num})

                # Extract images
                for img_idx, img_info in enumerate(page.get_images(full=True)):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        ext = base_image.get("ext", "png")
                        img_path = img_dir / f"page{page_num}_img{img_idx}.{ext}"
                        img_path.write_bytes(img_bytes)
                        images.append({
                            "type": "image",
                            "path": str(img_path),
                            "page": page_num,
                            "index": len(images),
                        })
                    except Exception as exc:
                        logger.warning(f"Failed to extract PDF image: {exc}")

            doc.close()

        except ImportError:
            logger.warning("PyMuPDF not installed; attempting pdf2image + OCR fallback")
            # Try pdf2image
            try:
                from pdf2image import convert_from_path

                pil_pages = convert_from_path(str(file_path), dpi=300)
                page_count = len(pil_pages)
                logger.info(f"pdf2image converted: {page_count} pages")

                out_dir = self._make_output_dir(file_path)
                img_dir = out_dir / "images"
                img_dir.mkdir(parents=True, exist_ok=True)

                for page_num, pil_img in enumerate(pil_pages):
                    img_path = img_dir / f"page_{page_num}.png"
                    pil_img.save(str(img_path))
                    images.append({
                        "type": "image",
                        "path": str(img_path),
                        "page": page_num,
                        "index": len(images),
                    })

                # OCR each page
                try:
                    ocr_engine = self._get_ocr_engine()
                    for page_num, pil_img in enumerate(pil_pages):
                        img_arr = np.array(pil_img)
                        result = ocr_engine.ocr(img_arr, cls=True)
                        if result and result[0]:
                            lines = [line[1][0] for line in result[0] if line[1]]
                            text = "\n".join(lines)
                            if text:
                                md_parts.append(text)
                                content_list.append({"type": "text", "text": text, "page": page_num})
                except ImportError:
                    logger.warning("PaddleOCR not available; OCR skipped for PDF fallback")

            except ImportError:
                logger.warning("pdf2image not installed; PDF fallback cannot proceed")

        return {
            "source_file": str(file_path),
            "pages": page_count,
            "content_list": content_list,
            "markdown": "\n\n".join(md_parts),
            "tables": tables,
            "images": images,
            "metadata": {
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser": "fallback_pdf",
            },
        }

    # ------------------------------------------------------------------
    # Cloud API integration (MinerU cloud service)
    # ------------------------------------------------------------------

    async def _parse_via_cloud_api(self, file_path: Path) -> dict | None:
        """
        Parse a document via MinerU cloud API when local models are unavailable.

        Supports two modes:
        - Agent lightweight API (no token, ≤10MB/≤20 pages, Markdown only)
        - Precise API (with token, ≤200MB/≤200 pages, Zip with MD+JSON)

        For PDFs exceeding 200 pages, automatically splits into ≤200-page chunks,
        processes each chunk via Precise API, then merges the results.

        Returns structured dict on success, None on failure.
        """
        import requests as sync_requests

        file_size = file_path.stat().st_size
        file_name = file_path.name
        suffix = _suffix(file_path)
        has_token = bool(self.api_token)

        # Determine which API to use
        use_precise = has_token
        use_agent = not has_token

        # Size checks
        if use_agent and file_size > 10 * 1024 * 1024:
            logger.warning(f"File too large for Agent API ({file_size/1024/1024:.1f}MB > 10MB limit)")
            use_agent = False
            if not has_token:
                logger.warning("No API token — cannot use Precise API either")
                return None

        if use_precise and file_size > 200 * 1024 * 1024:
            logger.error(f"File too large for any cloud API ({file_size/1024/1024:.1f}MB)")
            return None

        try:
            if use_precise:
                # Check page count for PDFs — split if > 200 pages
                if suffix == ".pdf":
                    page_count = self._count_pdf_pages(file_path)
                    if page_count > 200:
                        return await self._cloud_precise_split(file_path, file_name, page_count)

                return await self._cloud_precise_api(file_path, file_name, suffix)
            elif use_agent:
                return await self._cloud_agent_api(file_path, file_name, suffix)
            else:
                return None
        except Exception as exc:
            logger.error(f"MinerU cloud API failed: {exc}")
            return None

    @staticmethod
    def _count_pdf_pages(file_path: Path) -> int:
        """Count PDF pages using PyMuPDF (fitz)."""
        try:
            import fitz
            doc = fitz.open(str(file_path))
            n = len(doc)
            doc.close()
            return n
        except Exception:
            return 0

    def _split_pdf(self, file_path: Path, chunk_size: int = 200) -> list[Path]:
        """Split a PDF into chunks of ≤chunk_size pages, return temp file paths."""
        import fitz
        import tempfile

        doc = fitz.open(str(file_path))
        total = len(doc)
        chunks: list[Path] = []
        tmp_dir = Path(tempfile.mkdtemp(prefix="mineru_split_"))

        for start in range(0, total, chunk_size):
            end = min(start + chunk_size, total)
            chunk_doc = fitz.open()
            chunk_doc.insert_pdf(doc, from_page=start, to_page=end - 1)
            chunk_path = tmp_dir / f"{file_path.stem}_part{len(chunks)+1}_{start+1}-{end}.pdf"
            chunk_doc.save(str(chunk_path))
            chunk_doc.close()
            chunks.append(chunk_path)
            logger.info(f"Split chunk {len(chunks)}: pages {start+1}-{end} ({chunk_path.name})")

        doc.close()
        logger.info(f"Split {file_path.name} into {len(chunks)} chunks")
        return chunks

    async def _cloud_precise_split(self, file_path: Path, file_name: str, total_pages: int) -> dict | None:
        """Split a large PDF, process each chunk via Precise API, merge results."""
        import shutil

        logger.info(f"PDF has {total_pages} pages (>200), splitting into chunks")

        chunks = self._split_pdf(file_path, chunk_size=200)
        if not chunks:
            logger.error("PDF split produced no chunks")
            return None

        all_md_parts: list[str] = []
        all_content_list: list[dict] = []
        all_tables: list[dict] = []
        all_images: list[dict] = []
        total_pages_parsed = 0
        batch_ids: list[str] = []
        page_offset = 0

        try:
            for i, chunk_path in enumerate(chunks, 1):
                logger.info(f"Processing chunk {i}/{len(chunks)}: {chunk_path.name}")
                result = await self._cloud_precise_api(
                    chunk_path, chunk_path.name, ".pdf"
                )
                if result is None:
                    logger.warning(f"Chunk {i} failed, skipping")
                    continue

                md = result.get("markdown", "")
                if md:
                    all_md_parts.append(md)

                # Adjust page numbers in content_list
                for item in result.get("content_list", []):
                    if "page" in item:
                        item["page"] = item["page"] + page_offset
                    all_content_list.append(item)

                # Adjust table page references
                for tbl in result.get("tables", []):
                    if "page" in tbl:
                        tbl["page"] = tbl["page"] + page_offset
                    all_tables.append(tbl)

                all_images.extend(result.get("images", []))
                total_pages_parsed += result.get("pages", 0)
                page_offset += result.get("pages", 0)

                bid = result.get("metadata", {}).get("api_batch_id", "")
                if bid:
                    batch_ids.append(bid)

                logger.info(
                    f"Chunk {i} done: {result.get('pages', 0)} pages, "
                    f"{len(md)} chars, {len(result.get('tables', []))} tables"
                )
        finally:
            # Clean up temp split files
            try:
                shutil.rmtree(chunks[0].parent, ignore_errors=True)
            except Exception:
                pass

        if not all_md_parts:
            logger.error("All chunks failed — no results")
            return None

        merged_md = "\n\n---\n\n".join(all_md_parts)
        logger.info(
            f"Split processing complete: {total_pages_parsed}/{total_pages} pages, "
            f"{len(merged_md)} chars, {len(all_tables)} tables from {len(batch_ids)} batches"
        )

        return {
            "source_file": str(file_path),
            "pages": total_pages_parsed,
            "content_list": all_content_list,
            "markdown": merged_md,
            "tables": all_tables,
            "images": all_images,
            "metadata": {
                "file_name": file_name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_cloud_precise_split",
                "api_batch_ids": batch_ids,
                "total_chunks": len(chunks),
                "total_pages": total_pages,
            },
        }

    async def _cloud_agent_api(self, file_path: Path, file_name: str, suffix: str) -> dict | None:
        """Use the Agent lightweight API (no token, file upload mode)."""
        import requests as sync_requests
        import asyncio

        BASE_URL = "https://mineru.net/api/v1/agent"

        logger.info(f"MinerU cloud: using Agent lightweight API for {file_name}")

        # Step 1: Get signed upload URL
        data = {
            "file_name": file_name,
            "language": "ch",
            "enable_table": self.table_enable,
            "is_ocr": True,
            "enable_formula": self.formula_enable,
        }

        loop = asyncio.get_running_loop()
        resp = await loop.run_in_executor(
            None,
            lambda: sync_requests.post(f"{BASE_URL}/parse/file", json=data, timeout=30)
        )
        result = resp.json()

        if result.get("code") != 0:
            logger.error(f"Agent API submit failed: {result.get('msg', 'unknown')}")
            return None

        task_id = result["data"]["task_id"]
        file_url = result["data"]["file_url"]
        logger.info(f"Agent API task created: {task_id}")

        # Step 2: Upload file via PUT
        with open(file_path, "rb") as f:
            put_resp = await loop.run_in_executor(
                None,
                lambda: sync_requests.put(file_url, data=f, timeout=60)
            )
        if put_resp.status_code not in (200, 201):
            logger.error(f"File upload failed: HTTP {put_resp.status_code}")
            return None
        logger.info("File uploaded to cloud, waiting for parsing...")

        # Step 3: Poll for result
        markdown_content = await self._cloud_poll(
            f"{BASE_URL}/parse/{task_id}",
            timeout=300,
            extract_fn=lambda r: r["data"].get("markdown_url"),
            has_token=False,
        )

        if markdown_content is None:
            return None

        # Download the Markdown
        md_resp = await loop.run_in_executor(
            None,
            lambda: sync_requests.get(markdown_content, timeout=30)
        )
        md_text = md_resp.text

        logger.info(f"Agent API completed: {len(md_text)} chars Markdown")

        return {
            "source_file": str(file_path),
            "pages": 0,  # Unknown from Agent API
            "content_list": [{"type": "text", "text": md_text, "page": 0}],
            "markdown": md_text,
            "tables": [],  # Agent API doesn't return structured tables
            "images": [],
            "metadata": {
                "file_name": file_name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_cloud_agent",
                "api_task_id": task_id,
            },
        }

    async def _cloud_precise_api(self, file_path: Path, file_name: str, suffix: str) -> dict | None:
        """Use the Precise API (with token, file upload mode)."""
        import requests as sync_requests
        import asyncio
        import zipfile
        import io

        BASE_URL = "https://mineru.net/api/v4"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_token}",
        }

        logger.info(f"MinerU cloud: using Precise API for {file_name}")

        # Step 1: Get batch upload URL
        data = {
            "files": [{"name": file_name}],
            "model_version": "vlm",
            "enable_table": self.table_enable,
            "enable_formula": self.formula_enable,
        }

        loop = asyncio.get_running_loop()
        resp = await loop.run_in_executor(
            None,
            lambda: sync_requests.post(f"{BASE_URL}/file-urls/batch", json=data, headers=headers, timeout=30)
        )
        result = resp.json()

        if result.get("code") != 0:
            logger.error(f"Precise API submit failed: {result.get('msg', 'unknown')}")
            return None

        batch_id = result["data"]["batch_id"]
        upload_url = result["data"]["file_urls"][0]
        logger.info(f"Precise API batch created: {batch_id}")

        # Step 2: Upload file
        with open(file_path, "rb") as f:
            put_resp = await loop.run_in_executor(
                None,
                lambda: sync_requests.put(upload_url, data=f, timeout=120)
            )
        if put_resp.status_code not in (200, 201):
            logger.error(f"File upload failed: HTTP {put_resp.status_code}")
            return None
        logger.info("File uploaded, waiting for parsing...")

        # Step 3: Poll for result
        zip_url = await self._cloud_poll(
            f"{BASE_URL}/extract-results/batch/{batch_id}",
            timeout=600,
            extract_fn=lambda r: r["data"]["extract_result"][0].get("full_zip_url") if r["data"].get("extract_result") else None,
            has_token=True,
        )

        if zip_url is None:
            return None

        # Step 4: Download and parse the Zip
        zip_resp = await loop.run_in_executor(
            None,
            lambda: sync_requests.get(zip_url, timeout=60)
        )
        logger.info(f"Downloaded result zip: {len(zip_resp.content)} bytes")

        md_text = ""
        content_list = []
        tables = []
        page_count = 0

        with zipfile.ZipFile(io.BytesIO(zip_resp.content)) as zf:
            for name in zf.namelist():
                if name.endswith("full.md") or name.endswith(".md"):
                    md_text = zf.read(name).decode("utf-8")
                elif name.endswith("content_list.json"):
                    try:
                        cl = json.loads(zf.read(name).decode("utf-8"))
                        content_list = cl if isinstance(cl, list) else []
                    except Exception:
                        pass
                elif name.endswith("layout.json") or name.endswith("model.json"):
                    try:
                        layout = json.loads(zf.read(name).decode("utf-8"))
                        if isinstance(layout, list):
                            page_count = len(layout)
                    except Exception:
                        pass

        logger.info(f"Precise API completed: {len(md_text)} chars, {page_count} pages")

        return {
            "source_file": str(file_path),
            "pages": page_count,
            "content_list": content_list,
            "markdown": md_text,
            "tables": tables,
            "images": [],
            "metadata": {
                "file_name": file_name,
                "file_size": file_path.stat().st_size,
                "parser": "mineru_cloud_precise",
                "api_batch_id": batch_id,
            },
        }

    async def _cloud_poll(
        self,
        poll_url: str,
        timeout: int = 300,
        extract_fn = None,
        has_token: bool = False,
    ) -> str | None:
        """Poll a cloud API until done, return the extracted URL or None.

        Handles two response shapes:
        - Agent API:  state at ``data.state``
        - Batch API:  state at ``data.extract_result[0].state``
        """
        import requests as sync_requests
        import asyncio

        headers = {}
        if has_token:
            headers["Authorization"] = f"Bearer {self.api_token}"

        loop = asyncio.get_running_loop()
        t0 = time.perf_counter()
        interval = 5

        while time.perf_counter() - t0 < timeout:
            resp = await loop.run_in_executor(
                None,
                lambda: sync_requests.get(poll_url, headers=headers, timeout=30)
            )
            result = resp.json()
            data = result.get("data", {})

            # --- resolve state from either response shape ---
            state = data.get("state")
            if not state:
                # Batch (Precise) API: state is nested inside extract_result
                extract_result = data.get("extract_result")
                if isinstance(extract_result, list) and extract_result:
                    state = extract_result[0].get("state", "unknown")
                else:
                    state = "unknown"

            elapsed = int(time.perf_counter() - t0)

            if state == "done":
                url = extract_fn(result) if extract_fn else None
                if url:
                    logger.info(f"[{elapsed}s] Cloud parse completed")
                    return url
                logger.error("Cloud parse done but no result URL found")
                return None

            if state == "failed":
                err_msg = (
                    data.get("err_msg")
                    or data.get("extract_result", [{}])[0].get("err_msg", "unknown error")
                )
                logger.error(f"[{elapsed}s] Cloud parse failed: {err_msg}")
                return None

            # Still processing — try to show progress
            progress = data.get("extract_progress") or {}
            if progress:
                logger.info(
                    f"[{elapsed}s] {state}: "
                    f"{progress.get('extracted_pages', '?')}/{progress.get('total_pages', '?')} pages"
                )
            else:
                logger.info(f"[{elapsed}s] {state}...")

            await asyncio.sleep(interval)

        logger.error(f"Cloud poll timed out after {timeout}s")
        return None

    # ------------------------------------------------------------------
    # Generic fallback (last resort)
    # ------------------------------------------------------------------

    async def _fallback_parse(self, file_path: Path) -> dict:
        """Route to the appropriate fallback parser based on file extension."""
        suffix = _suffix(file_path)
        if suffix in _PDF_EXTS:
            return await self._parse_pdf_basic(file_path, preprocess=False)
        elif suffix in _IMAGE_EXTS:
            return await self._parse_image_basic(file_path, preprocess=False)
        elif suffix == ".docx":
            return await self._parse_docx_basic(file_path)
        elif suffix == ".pptx":
            return await self._parse_pptx_basic(file_path)
        elif suffix in _HTML_EXTS:
            return await self._parse_html(file_path)
        raise ValueError(f"No fallback parser for: {suffix}")

    # ------------------------------------------------------------------
    # Image pre-processing helpers
    # ------------------------------------------------------------------

    def _preprocess_pdf_pages(self, pdf_bytes: bytes) -> bytes:
        """Re-render PDF pages with image enhancement.

        Rasterises each page at 300 DPI, applies CLAHE + denoise + deskew,
        then re-encodes as a new PDF.
        """
        try:
            import fitz
        except ImportError:
            logger.warning("PyMuPDF not available; skipping PDF pre-processing")
            return pdf_bytes

        src_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        dst_doc = fitz.open()

        for page_idx in range(len(src_doc)):
            page = src_doc[page_idx]
            # Render at 300 DPI
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)

            enhanced = ImagePreprocessor.enhance(img)

            # Encode back to PNG
            pil_img = Image.fromarray(enhanced)
            buf = io.BytesIO()
            pil_img.save(buf, format="PNG")
            png_bytes = buf.getvalue()

            # Insert into new PDF page
            new_page = dst_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, stream=png_bytes)

        out_buf = io.BytesIO()
        dst_doc.save(out_buf)
        dst_doc.close()
        src_doc.close()
        return out_buf.getvalue()

    @staticmethod
    def _preprocess_image_bytes(img_bytes: bytes) -> bytes:
        """Enhance a single image and return the result as PNG bytes."""
        nparr = np.frombuffer(img_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        if img is None:
            return img_bytes
        enhanced = ImagePreprocessor.enhance(img)
        pil_img = Image.fromarray(enhanced)
        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Output directory management
    # ------------------------------------------------------------------

    def _make_output_dir(self, file_path: Path) -> Path:
        """Create a per-file output directory structure: <base>/<stem>/"""
        stem = file_path.stem
        # Sanitise for filesystem safety
        safe_stem = "".join(c if c.isalnum() or c in "-_." else "_" for c in stem)
        out_dir = self.output_base / safe_stem
        out_dir.mkdir(parents=True, exist_ok=True)
        return out_dir

    # ------------------------------------------------------------------
    # Markdown table helper
    # ------------------------------------------------------------------

    @staticmethod
    def _table_to_markdown(rows: list[list[str]]) -> str:
        """Convert a 2D list of strings into a Markdown table."""
        if not rows:
            return ""
        # Normalise column count
        max_cols = max(len(r) for r in rows)
        normalised = []
        for r in rows:
            padded = list(r) + [""] * (max_cols - len(r))
            normalised.append(padded)

        header = "| " + " | ".join(normalised[0]) + " |"
        separator = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
        body_lines = []
        for row in normalised[1:]:
            body_lines.append("| " + " | ".join(row) + " |")

        return "\n".join([header, separator] + body_lines)

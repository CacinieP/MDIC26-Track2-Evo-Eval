#!/usr/bin/env python3
"""
Build competition PPT for MDIC26 Track2 Evo-Eval
Usage: python scripts/build_ppt.py
Output: docs/presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# =============================================================================
# Color palette — Tech Blue (Radix-inspired)
# =============================================================================
C = {
    "dominant": RgbColor(0x1E, 0x66, 0xF5),
    "dominant_deep": RgbColor(0x1A, 0x5B, 0xD6),
    "secondary": RgbColor(0x6E, 0x56, 0xCF),
    "dark_bg": RgbColor(0x0B, 0x11, 0x20),
    "dark_bg_deep": RgbColor(0x06, 0x0A, 0x14),
    "white": RgbColor(0xFF, 0xFF, 0xFF),
    "ink": RgbColor(0x0C, 0x1A, 0x2B),
    "muted": RgbColor(0x5C, 0x6B, 0x7F),
    "light_bg": RgbColor(0xFB, 0xFC, 0xFF),
    "light_bg_alt": RgbColor(0xF0, 0xF4, 0xFF),
    "line": RgbColor(0xC6, 0xD2, 0xE0),
}

# Slide dimensions (16:9)
SW = Inches(13.333)
SH = Inches(7.5)

# Fonts
FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_CODE = "Consolas"


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, font_color,
                bold=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY):
    """Add a text box with consistent styling."""
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    """Add a shape with fill and optional line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size, font_color,
                    bullet_color=None, font_name=FONT_BODY):
    """Add a bulleted list."""
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.level = 0
        if bullet_color:
            p.font.color.rgb = bullet_color
    return tf


def footer(slide, n, total, dark=False):
    """Add footer with page number."""
    color = C["muted"] if not dark else RgbColor(0x88, 0x99, 0xAA)
    add_textbox(slide, Inches(0.5), SH - Inches(0.45), Inches(6), Inches(0.3),
                "Team CacinieP  |  MDIC26 Track2 Evo-Eval", Pt(10), color)
    add_textbox(slide, SW - Inches(1.5), SH - Inches(0.45), Inches(1), Inches(0.3),
                f"{n} / {total}", Pt(10), color, align=PP_ALIGN.RIGHT)


def stripe_bar(slide, dark=False):
    """Add top accent stripe."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.06), C["dominant"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(0.06), SW * 0.35, Inches(0.03), C["secondary"])


# =============================================================================
# Slide builders
# =============================================================================

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C["dark_bg"])

    # Decorative circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8),
              RgbColor(0x1E, 0x66, 0xF5))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(-1.8), Inches(-1.8), Inches(7.6), Inches(7.6), C["dark_bg"])

    # Kicker
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.5),
                "MDIC26 赛道二 · Data Agent 数据智能体", Pt(16), C["dominant"],
                bold=True, font_name=FONT_TITLE)

    # Main title
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2),
                "基于 MinerU 的智能文档处理系统", Pt(44), C["white"],
                bold=True, font_name=FONT_TITLE)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(3.7), Inches(10), Inches(0.6),
                "MinerU DataAgent — 自动化、高质量、低成本的复杂文档解析方案",
                Pt(18), RgbColor(0xD0, 0xD8, 0xE8), font_name=FONT_BODY)

    # Accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.5), Inches(1.5), Inches(0.06), C["dominant"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(4.5), Inches(0.5), Inches(0.06), C["secondary"])

    # Team info
    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.4),
                "团队：Team CacinieP", Pt(14), RgbColor(0x88, 0x99, 0xAA))
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(8), Inches(0.4),
                "技术栈：LangGraph · MinerU · FastAPI · Python 3.10+", Pt(12), RgbColor(0x66, 0x77, 0x88))

    footer(slide, 1, 12, dark=True)
    return slide


def build_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "PROJECT OVERVIEW", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "项目概述", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    # Three pillar cards
    cards = [
        ("任务一", "数据理解与结构化处理", "多源数据识别、信息抽取、格式转换、数据清洗、元数据抽取"),
        ("任务二", "复杂任务规划与自动执行", "LLM 自动拆解子任务、工具链调度、批处理、异常恢复"),
        ("任务三", "系统稳定性与综合能力", "高负载稳定执行、结果可靠性、325 自动化测试、Docker 部署"),
    ]
    for i, (kicker, title, body) in enumerate(cards):
        x = Inches(0.6 + i * 4.2)
        # Card bg
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.9), Inches(4.8),
                  C["light_bg_alt"], C["line"])
        # Kicker
        add_textbox(slide, x + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.3),
                    kicker, Pt(11), C["dominant"], bold=True)
        # Title
        add_textbox(slide, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.5),
                    title, Pt(16), C["ink"], bold=True)
        # Body
        add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(2.5),
                    body, Pt(12), C["muted"])

    footer(slide, 2, 12)
    return slide


def build_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "KEY CHALLENGES", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "五大难点场景与针对性方案", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    challenges = [
        ("财务报表数字解析", "表格结构识别 + LLM 校验 + 数值一致性检查（行合计/百分比/资产负债平衡）"),
        ("跨页合并与指代消解", "页面关联分析 + 实体注册表 + LLM 消解 + SequenceMatcher 续接检测"),
        ("低质量文档鲁棒性", "7 步自适应增强管线：CLAHE / 降噪 / 去倾斜 / 去印章 / 二值化 / DPI 提升"),
        ("复杂图表解析", "18 种图表分类 + LLM 视觉 + OpenCV 备选 + Chart-to-Table 转换"),
        ("国标/行业标准结构化", "MinerU 解析 + 后处理结构化 + 规范条款自动提取"),
    ]

    for i, (title, desc) in enumerate(challenges):
        y = Inches(1.5 + i * 1.1)
        # Number circle
        add_shape(slide, MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.4), Inches(0.4), C["dominant"])
        add_textbox(slide, Inches(0.6), y + Inches(0.05), Inches(0.4), Inches(0.3),
                    str(i + 1), Pt(14), C["white"], bold=True, align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, Inches(1.2), y, Inches(3.5), Inches(0.35),
                    title, Pt(14), C["ink"], bold=True)
        # Desc
        add_textbox(slide, Inches(1.2), y + Inches(0.3), Inches(11), Inches(0.5),
                    desc, Pt(12), C["muted"])

    footer(slide, 3, 12)
    return slide


def build_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["dark_bg"])

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "SYSTEM ARCHITECTURE", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "系统架构", Pt(28), C["white"], bold=True, font_name=FONT_TITLE)

    # Layer boxes
    layers = [
        ("API Gateway", "FastAPI REST API  |  CLI  |  Batch  |  Demo", Inches(1.0), C["dominant_deep"]),
        ("LangGraph Agent Core", "analyze_task → plan_execution → execute_step → verify_result → format_output",
         Inches(2.6), C["secondary"]),
        ("Tool Registry", "MinerU Parser  ·  Table Parser  ·  Chart Analyzer  ·  Image Enhancer  ·  Cross-Page Merger",
         Inches(4.2), RgbColor(0x3A, 0x4A, 0x5A)),
        ("Infrastructure", "MinerU Models (GPU/CPU)  ·  PaddleOCR  ·  OpenCV  ·  LLM API (Claude/Qwen/GLM)",
         Inches(5.8), RgbColor(0x2A, 0x3A, 0x4A)),
    ]

    for title, body, y, color in layers:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(1.2), color)
        add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(11), Inches(0.35),
                    title, Pt(16), C["white"], bold=True)
        add_textbox(slide, Inches(0.8), y + Inches(0.45), Inches(11), Inches(0.5),
                    body, Pt(12), RgbColor(0xCC, 0xDD, 0xEE))

    footer(slide, 4, 12, dark=True)
    return slide


def build_capability_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "CAPABILITY 01", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "复杂任务规划与自动执行", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    items = [
        "LangGraph StateGraph 6 节点 + 3 条件路由：自动拆解 → 有序执行 → 错误恢复 → 质量验证",
        "任务分析：LLM 深度分析（可用时）/ 启发式关键词匹配（降级时），自动识别文档类型与难度",
        "动态计划生成：基于 assessment 构建 SubTask DAG，自动添加表格提取、图表分析、跨页合并等步骤",
        "三级容错：单步 3 次重试 → 验证失败触发最多 2 轮重规划 → 最佳努力输出",
        "批处理支持：CLI batch 命令 + /tasks/batch API，支持 100+ 文件统一格式输出",
    ]
    for i, item in enumerate(items):
        y = Inches(1.6 + i * 0.95)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.08), Inches(0.55), C["dominant"])
        add_textbox(slide, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.5),
                    item, Pt(14), C["ink"])

    footer(slide, 5, 12)
    return slide


def build_capability_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "CAPABILITY 02", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "数据理解与结构化处理", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    # Two column layout
    # Left: tools
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(6), Inches(5.2),
              C["light_bg_alt"], C["line"])
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4),
                "核心工具链", Pt(16), C["ink"], bold=True)

    tools = [
        ("MinerU Parser", "PDF/图片/DOCX/PPTX/HTML 全格式支持，本地 + 云端双模"),
        ("Table Parser", "HTML→结构化网格 + 合并单元格展开 + 中文数字解析 + 数值自验证"),
        ("Chart Analyzer", "18 种图表分类 + LLM 视觉数据提取 + Chart-to-Table"),
        ("Image Enhancer", "7 步自适应管线：模糊/光照/签章/倾斜/降噪/二值化/超分"),
        ("Cross-Page Merger", "跨页表格续接检测 + 实体注册表 + 指代消解"),
    ]
    for i, (name, desc) in enumerate(tools):
        y = Inches(2.2 + i * 0.85)
        add_textbox(slide, Inches(0.85), y, Inches(5.4), Inches(0.3), name, Pt(13), C["dominant"], bold=True)
        add_textbox(slide, Inches(0.85), y + Inches(0.25), Inches(5.4), Inches(0.4), desc, Pt(11), C["muted"])

    # Right: output formats
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(5.2),
              C["light_bg_alt"], C["line"])
    add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5.5), Inches(0.4),
                "输出格式", Pt(16), C["ink"], bold=True)

    formats = [
        ("JSON", "完整结构化结果：execution_summary / tables / charts / verification / logs"),
        ("Markdown", "MinerU 原生 Markdown 输出，含表格和图片引用"),
        ("CSV", "表格数据可直接导出为 CSV 用于下游分析"),
    ]
    for i, (name, desc) in enumerate(formats):
        y = Inches(2.2 + i * 1.3)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), y, Inches(1.2), Inches(0.4), C["dominant"])
        add_textbox(slide, Inches(7.1), y + Inches(0.05), Inches(1.2), Inches(0.3),
                    name, Pt(11), C["white"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(8.4), y, Inches(4.3), Inches(0.6), desc, Pt(12), C["muted"])

    footer(slide, 6, 12)
    return slide


def build_capability_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "CAPABILITY 03", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "系统稳定性与综合能力评测", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    # Stats bar
    stats = [
        ("325", "自动化测试用例"),
        ("10", "核心测试模块"),
        ("0.7", "质量阈值"),
        ("3", "最大重试次数"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.6 + i * 3.2)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.9), Inches(1.4),
                  C["light_bg_alt"], C["line"])
        add_textbox(slide, x, Inches(1.6), Inches(2.9), Inches(0.7),
                    num, Pt(36), C["dominant"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_TITLE)
        add_textbox(slide, x, Inches(2.25), Inches(2.9), Inches(0.4),
                    label, Pt(12), C["muted"], align=PP_ALIGN.CENTER)

    # Details
    items = [
        "四维动态质量评分：内容完整性(0.30) + 执行完成度(0.30) + 任务匹配度(0.20) + 内容丰富度(0.20)",
        "异常恢复：工具未注册 → 跳过；MinerU 失败 → PyMuPDF 回退；内存溢出 → 自动分页",
        "工程化保障：Docker 多阶段构建 · GitHub Actions CI · Healthcheck · 结构化日志 + SSE 实时流",
        "双模部署：本地 GPU (25-30 页/秒) ↔ 云端 API (Precise/Free 自动 fallback)",
    ]
    for i, item in enumerate(items):
        y = Inches(3.2 + i * 0.95)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.08), Inches(0.55), C["dominant"])
        add_textbox(slide, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.5),
                    item, Pt(14), C["ink"])

    footer(slide, 7, 12)
    return slide


def build_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["dark_bg"])

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "TECHNICAL HIGHLIGHTS", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "技术创新点", Pt(28), C["white"], bold=True, font_name=FONT_TITLE)

    highlights = [
        ("中文大写数字解析器", "壹佰贰拾叁万 → 1,230,000，支持壹贰叁/一二三/万亿元单位自动转换"),
        ("数值自验证机制", "不仅提取数字，还能交叉验证：行合计核对、百分比验证、资产负债表平衡检查"),
        ("括号负数识别", "(1,234.56) → -1234.56，适配财务报告常见表示法"),
        ("大文件自动分片", "超过 200 页 PDF 自动切分 → 逐段解析 → 结果合并，实测 411 页报告 3 段处理"),
        ("无 LLM 降级运行", "未配置 API Key 时自动切换启发式分析，所有功能可用，精度略有下降"),
    ]

    for i, (title, desc) in enumerate(highlights):
        y = Inches(1.5 + i * 1.0)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.85),
                  RgbColor(0x1A, 0x2A, 0x3A))
        add_textbox(slide, Inches(0.85), y + Inches(0.08), Inches(3.5), Inches(0.3),
                    title, Pt(14), C["dominant"], bold=True)
        add_textbox(slide, Inches(4.5), y + Inches(0.08), Inches(7.8), Inches(0.5),
                    desc, Pt(12), RgbColor(0xCC, 0xDD, 0xEE))

    footer(slide, 8, 12, dark=True)
    return slide


def build_examples(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "TYPICAL EXAMPLES", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "典型任务执行示例与性能", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    # Table header
    headers = ["示例", "页数", "耗时", "输出", "质量分"]
    cols_x = [0.6, 4.5, 5.8, 7.3, 10.5]
    cols_w = [3.7, 1.1, 1.3, 3.0, 1.8]

    # Header bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12), Inches(0.45), C["dominant"])
    for i, h in enumerate(headers):
        add_textbox(slide, Inches(cols_x[i]), Inches(1.55), Inches(cols_w[i]), Inches(0.35),
                    h, Pt(12), C["white"], bold=True)

    rows = [
        ("金融公开披露材料", "411", "15.9s", "343,762 字符 Markdown + 28 表格", "0.85"),
        ("材料科学论文", "252", "~20s", "296,973 字符 Markdown + 图表", "0.85"),
        ("航天国标文件 (GB/T)", "8", "~5s", "4,472 字符 + 规范结构", "0.85"),
        ("自制财务样本", "1", "8.1s", "2 表格 + 数值一致性验证", "0.83"),
        ("411 页金融报告 (Cloud)", "411", "341.4s", "463,250 字符 (3 段合并)", "0.85"),
    ]

    for i, row in enumerate(rows):
        y = Inches(1.95 + i * 0.55)
        bg = C["light_bg_alt"] if i % 2 == 0 else C["white"]
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.5), bg)
        for j, cell in enumerate(row):
            add_textbox(slide, Inches(cols_x[j]), y + Inches(0.08), Inches(cols_w[j]), Inches(0.35),
                        cell, Pt(11), C["ink"])

    # Performance note
    add_textbox(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(0.4),
                "吞吐量：约 25–30 页/秒（GPU 本地模式）|  云端 Precise API：≤200MB/≤200 页/次",
                Pt(12), C["muted"])

    footer(slide, 9, 12)
    return slide


def build_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["white"])
    stripe_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "ENGINEERING & TESTING", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "工程化与测试覆盖", Pt(28), C["ink"], bold=True, font_name=FONT_TITLE)

    # Test modules grid
    modules = [
        ("test_api.py", "21", "TaskStore + API 端点"),
        ("test_graph.py", "27", "Agent 图/路由/验证"),
        ("test_mineru_parser.py", "38", "MinerU 解析/回退/分片"),
        ("test_chart_analyzer.py", "56", "图表分类/视觉/数据提取"),
        ("test_crosspage_merger.py", "50", "跨页合并/实体/指代消解"),
        ("test_image_enhancer.py", "58", "图像增强/质量评估/OCR 投票"),
    ]

    for i, (name, count, desc) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(1.5 + row * 2.5)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.2),
                  C["light_bg_alt"], C["line"])
        add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.6), Inches(0.3),
                    name, Pt(11), C["dominant"], bold=True, font_name=FONT_CODE)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.6), Inches(0.5),
                    f"{count} tests", Pt(24), C["dominant"], bold=True, font_name=FONT_TITLE)
        add_textbox(slide, x + Inches(0.15), y + Inches(1.15), Inches(3.6), Inches(0.5),
                    desc, Pt(11), C["muted"])

    footer(slide, 10, 12)
    return slide


def build_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["dark_bg"])

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                "APPLICATION VALUE", Pt(12), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
                "应用价值与适用场景", Pt(28), C["white"], bold=True, font_name=FONT_TITLE)

    scenes = [
        ("金融", "年报/季报批量解析、财务数据提取", "提取效率提升 10–50 倍"),
        ("法律", "合同文档结构化、法规条款提取", "降低人力成本 80%"),
        ("制造/航天", "技术文档解析、国标规范提取", "自动化合规检查"),
        ("学术研究", "论文批量处理、图表数据提取", "高质量训练语料生产"),
        ("教育", "教材数字化、图文混合处理", "知识图谱构建"),
    ]

    for i, (industry, scene, value) in enumerate(scenes):
        y = Inches(1.5 + i * 1.0)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.85),
                  RgbColor(0x1A, 0x2A, 0x3A))
        # Industry tag
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), y + Inches(0.2), Inches(1.0), Inches(0.35),
                  C["dominant"])
        add_textbox(slide, Inches(0.85), y + Inches(0.22), Inches(1.0), Inches(0.3),
                    industry, Pt(11), C["white"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.1), y + Inches(0.08), Inches(5.5), Inches(0.3),
                    scene, Pt(13), C["white"])
        add_textbox(slide, Inches(7.8), y + Inches(0.08), Inches(4.5), Inches(0.3),
                    value, Pt(12), C["dominant"], bold=True, align=PP_ALIGN.RIGHT)

    footer(slide, 11, 12, dark=True)
    return slide


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C["dark_bg"])

    # Decorative circles
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6), RgbColor(0x1E, 0x66, 0xF5))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.2), Inches(-1.8), Inches(5.6), Inches(5.6), C["dark_bg"])

    add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5),
                "THANK YOU", Pt(16), C["dominant"], bold=True)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(10), Inches(1.0),
                "基于 MinerU 的 Data Agent", Pt(44), C["white"], bold=True, font_name=FONT_TITLE)
    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(10), Inches(0.5),
                "自动化 · 高质量 · 低成本 · 可复现", Pt(20), RgbColor(0xCC, 0xDD, 0xEE))

    # QR hint
    add_textbox(slide, Inches(0.6), Inches(4.5), Inches(10), Inches(0.3),
                "GitHub: github.com/CacinieP/MDIC26-Track2-Evo-Eval", Pt(12), C["muted"])

    # Accent bars
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.2), Inches(1.5), Inches(0.06), C["dominant"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(5.2), Inches(0.5), Inches(0.06), C["secondary"])

    footer(slide, 12, 12, dark=True)
    return slide


# =============================================================================
# Main
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    build_cover(prs)
    build_overview(prs)
    build_challenges(prs)
    build_architecture(prs)
    build_capability_1(prs)
    build_capability_2(prs)
    build_capability_3(prs)
    build_highlights(prs)
    build_examples(prs)
    build_engineering(prs)
    build_value(prs)
    build_closing(prs)

    output_dir = Path(__file__).resolve().parent.parent / "docs"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "presentation.pptx"
    prs.save(str(output_path))
    print(f"[OK] Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
